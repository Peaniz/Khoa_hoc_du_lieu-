from pptx import Presentation
import os

def read_pptx(file_path):
    # Load the presentation
    prs = Presentation(file_path)
    
    print("=== Nội dung bài giảng ===\n")
    
    # Iterate through all slides
    for slide_number, slide in enumerate(prs.slides, 1):
        print(f"\nSlide {slide_number}:")
        print("-" * 50)
        
        # Get slide title if it exists
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    title = shape.text
                    print(f"Tiêu đề: {title}")
                else:
                    print(f"Nội dung: {shape.text}")
        
        print("-" * 50)

# Read the PowerPoint file
pptx_path = "Chap 2. Các công cụ toán học trong khoa học dữ liệu - phan 1_.pptx"
read_pptx(pptx_path) 